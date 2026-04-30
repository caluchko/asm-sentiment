"""
Generate a presentation-ready .pptx file for the ASM media sentiment analysis.
Run from the project root: python3 scripts/create_deck.py
Output: exports/asm_sentiment_deck.pptx  (open in Keynote or PowerPoint)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Colours (match chart palette)
# ---------------------------------------------------------------------------
NAVY      = RGBColor(0x1F, 0x4E, 0x79)   # chart rolling-average blue
MID_BLUE  = RGBColor(0x2E, 0x75, 0xB6)   # accent
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)   # body text
MID_GRAY  = RGBColor(0x77, 0x77, 0x77)   # captions / secondary
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# ---------------------------------------------------------------------------
# Slide geometry  (16 : 9 widescreen)
# ---------------------------------------------------------------------------
W = Inches(13.333)
H = Inches(7.5)

MARGIN_L = Inches(0.55)
MARGIN_T = Inches(0.40)
CONTENT_W = Inches(12.23)

OUTPUT = os.path.join("exports", "asm_sentiment_deck.pptx")
EXPORT_DIR = "exports"

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank


def add_title(slide, text: str,
              top=Inches(0.35), height=Inches(0.75),
              size=Pt(34), color=NAVY, bold=True):
    txb = slide.shapes.add_textbox(MARGIN_L, top, CONTENT_W, height)
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_rule(slide, top=Inches(1.18), color=MID_BLUE):
    """Thin horizontal rule under the title."""
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE (freeform not needed for a thin rect)
        MARGIN_L, top, CONTENT_W, Inches(0.025),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_bullets(slide, items: list[str | tuple],
                top=Inches(1.35), height=Inches(5.85),
                base_size=Pt(21), indent_size=Pt(19)):
    """
    items: list of str (top-level) or (str, level) tuples where level 1 = sub-bullet.
    """
    txb = slide.shapes.add_textbox(MARGIN_L, top, CONTENT_W, height)
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False

        p.level = level
        p.space_before = Pt(4) if level == 1 else Pt(10)

        run = p.add_run()
        run.text = ("    • " if level == 1 else "• ") + text
        run.font.size = indent_size if level == 1 else base_size
        run.font.color.rgb = MID_GRAY if level == 1 else DARK_GRAY


def add_chart_image(slide, img_path: str,
                    top=Inches(1.25), padding_lr=Inches(0.3)):
    """Insert a chart PNG, scaled to fill the slide below the title."""
    avail_w = W - 2 * padding_lr
    avail_h = H - top - Inches(0.15)
    slide.shapes.add_picture(img_path, padding_lr, top, avail_w, avail_h)


def add_caption(slide, text: str, top=None):
    """Small italic caption near the bottom of the slide."""
    cap_h = Inches(0.35)
    t = top if top is not None else H - cap_h - Inches(0.1)
    txb = slide.shapes.add_textbox(MARGIN_L, t, CONTENT_W, cap_h)
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.italic = True
    run.font.color.rgb = MID_GRAY


# ---------------------------------------------------------------------------
# Build deck
# ---------------------------------------------------------------------------

prs = new_deck()

# ── Slide 1: Title ──────────────────────────────────────────────────────────
s = blank_slide(prs)

# Main title
txb = slide = s  # reuse variable name for brevity
tb = s.shapes.add_textbox(MARGIN_L, Inches(2.2), CONTENT_W, Inches(1.4))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Artisanal and Small-Scale Mining"
run.font.size = Pt(44)
run.font.bold = True
run.font.color.rgb = NAVY

# Subtitle line 1
p2 = tf.add_paragraph()
run2 = p2.add_run()
run2.text = "Media Sentiment Analysis"
run2.font.size = Pt(44)
run2.font.bold = True
run2.font.color.rgb = NAVY

# Subtitle line 2
p3 = tf.add_paragraph()
p3.space_before = Pt(16)
run3 = p3.add_run()
run3.text = "GDELT DOC 2.0 API  ·  2017–present"
run3.font.size = Pt(22)
run3.font.bold = False
run3.font.color.rgb = MID_BLUE

# Rule beneath title block
add_rule(s, top=Inches(4.0))

# Footer
add_caption(s, "Data source: GDELT Project (gdeltproject.org)  ·  WB_555 Artisanal and Small-Scale Mining theme")

# ── Slide 2: Introduction & Rationale ───────────────────────────────────────
s = blank_slide(prs)
add_title(s, "Introduction & Rationale")
add_rule(s)
add_bullets(s, [
    "Artisanal and small-scale mining (ASM) employs an estimated 40–100 million people worldwide, "
    "yet remains poorly understood by mainstream policy and finance",
    "Media coverage shapes how ASM is perceived by governments, investors, and the public — "
    "influencing regulation, funding, and support for the sector",
    "Tracking sentiment over time reveals whether discourse is improving, worsening, "
    "or shifting in response to key events and policy milestones",
    "This analysis provides a systematic, data-driven baseline for understanding "
    "how global media has framed ASM from 2017 to the present",
])

# ── Slide 3: Methodology ────────────────────────────────────────────────────
s = blank_slide(prs)
add_title(s, "Methodology")
add_rule(s)
add_bullets(s, [
    "Data source: GDELT DOC 2.0 API — a real-time index of news media across hundreds of thousands "
    "of outlets in 65+ languages",
    "Filter: GDELT GKG theme tag WB_555 (World Bank taxonomy: Artisanal and Small-Scale Mining), "
    "applied automatically by GDELT's entity-recognition system",
    "Date range: January 2017 to present (~3 385 daily observations)",
    "Sentiment metric: GDELT average tone — a composite score from 40+ sentiment dictionaries "
    "applied to full article text; roughly −10 (most negative) to +10 (most positive)",
    "Supplementary queries: seven keyword searches (galamsey, garimpeiro, orpaillage, etc.) "
    "and four comparison themes (large-scale mining, gold, extractives, deforestation)",
    "Smoothing: 30-day rolling average applied to reduce day-to-day noise",
])

# ── Slide 4: Chart — Sentiment Trend ────────────────────────────────────────
s = blank_slide(prs)
add_title(s, "Sentiment Trend (2017–present)", height=Inches(0.6), size=Pt(26))
add_chart_image(s, "exports/01_sentiment_trend.png", top=Inches(1.0))

# ── Slide 5: Chart — Coverage Volume ────────────────────────────────────────
s = blank_slide(prs)
add_title(s, "Coverage Volume", height=Inches(0.6), size=Pt(26))
add_chart_image(s, "exports/02_coverage_volume.png", top=Inches(1.0))

# ── Slide 6: Chart — Coverage by Country ────────────────────────────────────
s = blank_slide(prs)
add_title(s, "Coverage by Source Country (top 8, 30-day rolling)", height=Inches(0.6), size=Pt(26))
add_chart_image(s, "exports/03_coverage_by_country.png", top=Inches(1.0))

# ── Slide 7: Chart — Keyword Sentiment ──────────────────────────────────────
s = blank_slide(prs)
add_title(s, "Sentiment by Keyword", height=Inches(0.6), size=Pt(26))
add_chart_image(s, "exports/04_keyword_sentiment.png", top=Inches(1.0))

# ── Slide 8: Chart — Annual Context ─────────────────────────────────────────
s = blank_slide(prs)
add_title(s, "Annual Sentiment in Context", height=Inches(0.6), size=Pt(26))
add_chart_image(s, "exports/05_annual_context.png", top=Inches(1.0))

# ── Slide 9: Key Findings ───────────────────────────────────────────────────
s = blank_slide(prs)
add_title(s, "Key Findings")
add_rule(s)
add_bullets(s, [
    "ASM sentiment is persistently negative — mean tone ~−0.67 across 2017–present; "
    "60% of days are below zero",
    "Modest long-run improvement: mean tone −0.86 in 2017–18 vs. −0.65 in 2024–25, "
    "suggesting gradually less negative framing",
    "2023 anomaly: the least negative year on record (median ≈ 0); "
    "driven partly by increased positive investment and reform coverage",
    "Ghana dominates coverage by volume intensity; galamsey-related articles "
    "are consistently the most negative keyword series",
    "ASM is framed more negatively than large-scale mining and the extractives sector broadly, "
    "but less negatively than deforestation",
    "Coverage spikes align with major events: Ghana crackdowns (2021), "
    "recent Zambia and Zimbabwe activity (2025–26)",
])

# ── Slide 10: Unknowns & Caveats ────────────────────────────────────────────
s = blank_slide(prs)
add_title(s, "Unknowns & Caveats")
add_rule(s)
add_bullets(s, [
    "WB_555 false-positive rate is ~40% — many articles are off-topic "
    "(lithium, silver, unrelated mining); noise appears stable over time so trend comparisons remain valid",
    "Sentiment is a daily aggregate — GDELT DOC API does not expose per-article scores; "
    "within-day variance, medians, and outliers are not available at this level",
    "Source country reflects publication registration, not the story's subject country — "
    "e.g. allafrica.com articles appear as Nigeria regardless of coverage geography",
    "GDELT's source base has expanded significantly since 2017; "
    "raw article counts are normalised per 100k corpus to correct for this",
    "Non-English articles are machine-translated before sentiment scoring, "
    "which may introduce noise for French (orpaillage) and Portuguese (garimpeiro) sources",
])

# ── Slide 11: Next Steps — Phase 2 ──────────────────────────────────────────
s = blank_slide(prs)
add_title(s, "Possible Next Steps — Phase 2")
add_rule(s)
add_bullets(s, [
    "Per-article tone via BigQuery GKG — access individual article scores to compute "
    "daily medians, standard deviations, and identify the most extreme articles",
    "Geographic disambiguation — link coverage to the country being reported on "
    "(not just the outlet's registration country) using GKG location fields",
    "Entity and event analysis — identify which actors, organisations, and event types "
    "drive the most negative and most positive coverage",
    "Correlate sentiment with external data — overlay production volumes, gold price, "
    "regulatory events, and development finance flows",
    "Expand keyword coverage — add country-specific terms and languages "
    "to reduce gaps in non-English ASM coverage",
])

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
os.makedirs(EXPORT_DIR, exist_ok=True)
prs.save(OUTPUT)
print(f"Saved → {OUTPUT}")
